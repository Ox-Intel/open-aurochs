import os
from io import BytesIO
from PIL import Image, ImageOps, ImageDraw, ImageFont
from django.conf import settings
from pptx import Presentation


EXPORT_ASSETS_DIR = os.path.join(
    settings.APPS_DIR, "api", "chart_utils", "export_assets"
)


def _add_header_footer(
    image_bytes: BytesIO, include_title: bool, title: str = None
) -> BytesIO:
    """
    Helper function that adds the header and footer to an image. Adds the ox logo to the bottom
    and optionally a title to the top and returns a bytesIo object with the image.
    Args:
        image_bytes: bytesIo file like object with the image data
        include_title: boolean to include a tile in the image
        title: Title string to be added to the image

    Returns:
        BytesIO file like object with image data.

    """

    # todo need a tmp dir to write to @steven where?
    font_path = os.path.join(EXPORT_ASSETS_DIR, "FjallaOne-Regular.ttf")
    logo_path = os.path.join(EXPORT_ASSETS_DIR, "ox-logo.jp2")
    logo = Image.open(logo_path).convert("RGBA").resize((45, 45))

    img = Image.open(image_bytes)
    if include_title:
        img = ImageOps.pad(
            img, size=(img.width + 40, img.height + 65), centering=(0.5, 1)
        )
        draw = ImageDraw.Draw(img)
        font = ImageFont.truetype(font_path, 28)
        draw.text((30, 20), title, (119, 119, 119), font=font)
        draw.line(
            [(20, 55), (img.width - 30, 55)], fill=(119, 119, 119), width=2, joint=None
        )

    # add logo to lower right corner
    img.paste(logo, (img.width - 90, img.height - 70))

    background = Image.new("RGBA", img.size, (255, 255, 255))
    output_img = Image.alpha_composite(background, img)

    output = BytesIO()

    output_img.save(output, format="PNG")

    return output


def assemble_ppt(image_list: list) -> Presentation:
    """
    Takes a list of image image and a list of titles and creates slides for each image. Returns a ppt file.
    Args:
        image_list: list of dicts with image key bytesIo objects with image data and title key for the slide title.

    Returns: Presentation

    """

    # Template file
    template_file = f"{EXPORT_ASSETS_DIR}/ppt_template.pptx"

    prs = Presentation(template_file)

    chart_slide_layout = prs.slide_layouts[0]

    for item in image_list:
        slide = prs.slides.add_slide(chart_slide_layout)

        header_title = slide.shapes.title
        header_title.text = item.get("title")

        chart_placeholder = slide.placeholders[10]

        chart_placeholder.insert_picture(item.get("image"))

    return prs
